"""
app/services/file_loader.py
Loads PDF, CSV, Excel, DOCX, PPTX, TXT, MD, JSON → list of chunk dicts.

Each chunk dict:
  {"text": str, **metadata_fields}
  e.g. {"text": "...", "page": 3}
       {"text": "...", "sheet": "Employees", "row": 12}

Key fixes vs original:
  - dict merge was broken: {"text": record, *meta} → now uses {**meta, "text": record}
  - header_summary chunk dict was malformed
  - Added PPTX, JSON loaders
  - Proper Excel hierarchy detection (builds org chart text block)
  - chardet for CSV encoding detection
  - Correct use of os.path vs Path
"""
from __future__ import annotations

import json
import os
import uuid
from pathlib import Path
from typing import Dict, List, Tuple

import chardet
import pandas as pd
from pypdf import PdfReader

from app.config import get_settings
from app.utils.logger import logger

settings = get_settings()

# ── Text splitter ─────────────────────────────────────────────────────────────

def _chunk_text(text: str, chunk_size: int | None = None, overlap: int | None = None) -> List[str]:
    """Word-based sliding window splitter."""
    chunk_size = chunk_size or settings.CHUNK_SIZE
    overlap = overlap or settings.CHUNK_OVERLAP
    if not text or not text.strip():
        return []
    words = text.split()
    chunks: List[str] = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i : i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        i += chunk_size - overlap
    return chunks


# ── Per-format loaders ────────────────────────────────────────────────────────

def _load_pdf(path: str) -> List[Dict]:
    reader = PdfReader(path)
    chunks: List[Dict] = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        for c in _chunk_text(text):
            chunks.append({"text": c, "page": page_num})
    return chunks


def _load_txt(path: str) -> List[Dict]:
    raw = Path(path).read_bytes()
    enc = chardet.detect(raw)["encoding"] or "utf-8"
    text = raw.decode(enc, errors="ignore")
    return [{"text": c} for c in _chunk_text(text)]


def _load_json(path: str) -> List[Dict]:
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    text = json.dumps(data, indent=2)
    return [{"text": c} for c in _chunk_text(text)]


def _load_docx(path: str) -> List[Dict]:
    from docx import Document  # python-docx
    doc = Document(path)
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": c} for c in _chunk_text(full_text)]


def _load_pptx(path: str) -> List[Dict]:
    from pptx import Presentation
    prs = Presentation(path)
    chunks: List[Dict] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            combined = "\n".join(texts)
            for c in _chunk_text(combined):
                chunks.append({"text": c, "slide": slide_num})
    return chunks


def _load_csv(path: str) -> List[Dict]:
    raw = Path(path).read_bytes()
    enc = chardet.detect(raw)["encoding"] or "utf-8"
    df = pd.read_csv(path, encoding=enc, dtype=str).fillna("")
    return _df_to_chunks(df)


def _load_excel(path: str) -> List[Dict]:
    xl = pd.ExcelFile(path)
    chunks: List[Dict] = []
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name, dtype=str).fillna("")
        for chunk in _df_to_chunks(df, sheet=sheet_name):
            chunks.append(chunk)
    return chunks


# ── DataFrame → chunks ────────────────────────────────────────────────────────

# Column names that indicate a reporting/manager relationship
_HIERARCHY_KEYWORDS = {
    "reports to", "reports_to", "manager", "reporting manager",
    "line manager", "managed by", "supervisor", "parent",
    "reporting to", "head",
}

_NAME_KEYWORDS = {
    "name", "employee name", "employee", "full name",
    "person", "staff", "staff name",
}


def _df_to_chunks(df: pd.DataFrame, sheet: str | None = None) -> List[Dict]:
    """
    Convert a DataFrame to searchable text chunks.
    1. Column header summary
    2. Hierarchy text block (if reporting column detected)
    3. One chunk per row (natural-language key:value format)
    """
    if df.empty:
        return []

    df.columns = [str(c).strip() for c in df.columns]
    col_lower = {c.lower(): c for c in df.columns}
    base_meta: Dict = {}
    if sheet:
        base_meta["sheet"] = sheet

    chunks: List[Dict] = []

    # 1. Header summary
    header_text = f"{'Sheet: ' + sheet + ' | ' if sheet else ''}Columns: {', '.join(df.columns)} | Total rows: {len(df)}"
    chunks.append({**base_meta, "text": header_text, "row": 0})

    # 2. Hierarchy block (if applicable)
    hier_chunk = _build_hierarchy_chunk(df, col_lower, sheet)
    if hier_chunk:
        chunks.append(hier_chunk)

    # 3. Row-level chunks
    for idx, row in df.iterrows():
        parts = [
            f"{col}: {str(row[col]).strip()}"
            for col in df.columns
            if str(row[col]).strip()
        ]
        if not parts:
            continue
        record = " | ".join(parts)
        chunks.append({**base_meta, "text": record, "row": int(idx) + 1})

    return chunks


def _build_hierarchy_chunk(
    df: pd.DataFrame, col_lower: Dict[str, str], sheet: str | None
) -> Dict | None:
    """
    If the DataFrame has a manager/reports-to column, build a full
    "X reports to Y" paragraph so the LLM can answer hierarchy questions
    without needing to piece it together from row chunks.
    """
    manager_col = next(
        (col_lower[kw] for kw in _HIERARCHY_KEYWORDS if kw in col_lower), None
    )
    if not manager_col:
        return None

    name_col = next(
        (col_lower[kw] for kw in _NAME_KEYWORDS if kw in col_lower), None
    )
    if not name_col:
        return None

    lines = [f"Reporting hierarchy{' in sheet ' + sheet if sheet else ''}:"]
    for _, row in df.iterrows():
        emp = str(row.get(name_col, "")).strip()
        mgr = str(row.get(manager_col, "")).strip()
        if emp and mgr and emp.lower() not in ("nan", "") and mgr.lower() not in ("nan", ""):
            lines.append(f"  {emp} reports to {mgr}")

    if len(lines) <= 1:
        return None

    meta: Dict = {"type": "hierarchy"}
    if sheet:
        meta["sheet"] = sheet
    meta["text"] = "\n".join(lines)
    return meta


# ── Dispatcher ────────────────────────────────────────────────────────────────

LOADERS: Dict[str, callable] = {
    ".pdf":  _load_pdf,
    ".csv":  _load_csv,
    ".xlsx": _load_excel,
    ".xls":  _load_excel,
    ".xlsm": _load_excel,
    ".docx": _load_docx,
    ".doc":  _load_docx,
    ".pptx": _load_pptx,
    ".ppt":  _load_pptx,
    ".txt":  _load_txt,
    ".md":   _load_txt,
    ".json": _load_json,
}

SUPPORTED_EXTENSIONS = set(LOADERS.keys())


def extract_chunks(filepath: str) -> List[Dict]:
    ext = Path(filepath).suffix.lower()
    loader = LOADERS.get(ext)
    if not loader:
        raise ValueError(f"Unsupported file type '{ext}'. Supported: {sorted(SUPPORTED_EXTENSIONS)}")
    logger.info(f"Loading [{ext}] {filepath}")
    chunks = loader(filepath)
    logger.info(f"  → {len(chunks)} chunks extracted")
    return chunks


def save_upload(file_bytes: bytes, original_name: str) -> Tuple[str, str]:
    """Save raw bytes to uploads dir. Returns (file_id, full_path)."""
    Path(settings.UPLOAD_DIR).mkdir(parents=True, exist_ok=True)
    file_id = str(uuid.uuid4())
    ext = Path(original_name).suffix.lower()
    dest = os.path.join(settings.UPLOAD_DIR, f"{file_id}{ext}")
    Path(dest).write_bytes(file_bytes)
    return file_id, dest
